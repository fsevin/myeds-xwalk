"""Generate AI Image (ai-image) block implementation guide PPTX (Adobe branding)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Adobe brand colours
RED = RGBColor(0xFF, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LGREY = RGBColor(0xF0, 0xF0, 0xF0)
DGREY = RGBColor(0x33, 0x33, 0x33)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


# ── helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_bullet_slide(title_text, bullets, code=None):
    slide = prs.slides.add_slide(BLANK)

    add_rect(slide, 0, 0, W, H, WHITE)
    add_rect(slide, 0, 0, Inches(0.18), H, RED)

    add_text(slide, title_text,
             Inches(0.35), Inches(0.3), Inches(12.5), Inches(0.8),
             size=28, bold=True, color=RED)

    sep = slide.shapes.add_shape(1, Inches(0.35), Inches(1.05), Inches(12.6), Pt(2))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RED
    sep.line.fill.background()

    y = Inches(1.25)
    for bullet in bullets:
        indent = bullet.startswith("  ")
        txt = bullet.lstrip()
        bsize = 16 if not indent else 14
        bcolor = DGREY
        prefix = "• " if not indent else "    – "
        add_text(slide, prefix + txt,
                 Inches(0.5), y, Inches(12.2), Inches(0.45),
                 size=bsize, color=bcolor)
        y += Inches(0.42 if not indent else 0.38)

    if code:
        cy = y + Inches(0.1)
        ch = H - cy - Inches(0.3)
        add_rect(slide, Inches(0.5), cy, Inches(12.3), ch, CODE_BG)
        txb = slide.shapes.add_textbox(Inches(0.65), cy + Inches(0.1),
                                        Inches(12.0), ch - Inches(0.2))
        txb.word_wrap = False
        tf = txb.text_frame
        tf.word_wrap = False
        first = True
        for line in code.split("\n"):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.size = Pt(11.5)
            run.font.bold = False
            run.font.color.rgb = DGREY
            run.font.name = "Courier New"

    add_text(slide, "Adobe",
             Inches(0.35), H - Inches(0.5), Inches(1.5), Inches(0.4),
             size=13, bold=True, color=RED)
    add_text(slide, "©2026 Adobe. All Rights Reserved. Adobe Confidential.",
             Inches(8), H - Inches(0.5), Inches(5), Inches(0.4),
             size=9, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)

    return slide


# ── slide 1 — title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, RED)
add_text(s, "The ai-image Block",
         Inches(0.6), Inches(2.8), Inches(9), Inches(1),
         size=40, bold=False, color=WHITE)
add_text(s, "Author-Generated, DAM-Persisted Firefly Images in Universal Editor",
         Inches(0.6), Inches(3.75), Inches(11), Inches(0.6),
         size=20, bold=True, color=WHITE)
add_text(s, "Adobe Experience Manager — Developer Guide",
         Inches(0.6), Inches(4.35), Inches(8), Inches(0.5),
         size=18, bold=False, color=WHITE)
add_text(s, "Adobe",
         Inches(0.6), Inches(1.8), Inches(2), Inches(0.6),
         size=22, bold=True, color=WHITE)

# ── slide 2 — agenda ─────────────────────────────────────────────────────────
add_bullet_slide("Agenda", [
    "Why ai-image is a separate block from firefly",
    "Architecture overview",
    "Step 1 — Authoring model: prompt, aspect ratio, persisted image",
    "Step 2 — Edge Worker route that streams image bytes",
    "Step 3 — Detecting Universal Editor authoring context",
    "Step 4 — Auto-generate on prompt entry (debounced)",
    "Step 5 — Upload to AEM DAM (direct binary upload API)",
    "Step 6 — Patch the block's own content node",
    "Bugs hit along the way — and how each was diagnosed",
    "Current status & open issue",
    "Best Practices & Summary",
])

# ── slide 3 — why a separate block ───────────────────────────────────────────
add_bullet_slide("Why a Separate Block From firefly", [
    "The original firefly block calls Firefly on every single render —",
    "author preview and every live visitor alike",
    "  Costs a Firefly credit per page view",
    "  Different image every time (no determinism)",
    "  The returned URL is a presigned S3 link that expires in ~1 hour",
    "",
    "ai-image fixes all three by generating once, in Universal Editor,",
    "and persisting the result as a real AEM DAM asset",
    "— published pages then render that stored asset with zero API calls",
    "",
    "Kept as a new block folder so the original firefly block stays untouched",
])

# ── slide 4 — architecture ────────────────────────────────────────────────────
add_bullet_slide("Architecture Overview", [
    "1. Author drops the block and types a Prompt (+ optional aspect ratio)",
    "2. 1.2s after the last edit, the block calls the edge Worker for image bytes",
    "3. Bytes are shown immediately as an inline preview",
    "4. Block fetches a Granite CSRF token, then uploads the bytes to AEM DAM",
    "   via the 3-step direct binary upload API",
    "5. Block POSTs the resulting DAM path as its own \"image\" property",
    "   (same-origin Sling POST against its own content node)",
    "6. Page reloads — Universal Editor now sees a real persisted asset",
    "7. On any later render (author or published), decorate() finds the",
    "   <picture> already present and returns immediately — no API calls",
], code="""\
Author types prompt
        │  (debounced 1.2s)
        ▼
ai-image.js  ──POST /api/firefly/generate-image──▶  edge/api-proxy.js (Worker)
        │  image bytes                                     │
        ▼                                     Adobe IMS + Firefly Generate API
  inline <img> preview
        │
        ├──GET /libs/granite/csrf/token.json (same-origin, cookies)
        ├──POST .initiateUpload.json / PUT blob / POST .completeUpload.json
        └──POST <block's own resource path>  { image: <dam-path> }
                        │
                 window.location.reload()""")

# ── slide 5 — step 1: model ───────────────────────────────────────────────────
add_bullet_slide("Step 1 — Authoring Model", [
    "Same prompt/aspectRatio shape as the firefly block, plus a third field:",
    "\"image\" — a reference (asset picker) field, written to programmatically",
    "instead of through the picker UI, but still manually overridable there",
], code="""\
// blocks/ai-image/_ai-image.json
{
  "models": [{
    "id": "ai-image",
    "fields": [
      { "component": "richtext", "name": "prompt", "label": "Prompt",
        "valueType": "string", "required": true },
      { "component": "select", "name": "aspectRatio", "label": "Aspect ratio",
        "valueType": "string",
        "options": [
          { "name": "Square (1024x1024)",   "value": "1024x1024" },
          { "name": "Landscape (1344x768)", "value": "1344x768" },
          { "name": "Portrait (768x1344)",  "value": "768x1344" }
        ] },
      { "component": "reference", "name": "image", "label": "Generated image",
        "valueType": "string" }
    ]
  }]
}""")

# ── slide 6 — step 2: worker route ───────────────────────────────────────────
add_bullet_slide("Step 2 — Edge Worker: Stream Image Bytes", [
    "Reuses the same IMS client_credentials auth as the firefly block",
    "(getFireflyToken, cached in-isolate) — no duplicated auth logic",
    "But instead of returning the presigned Firefly URL as JSON, it fetches",
    "the bytes server-side and streams them straight back",
    "Avoids S3 CORS uncertainty and the URL's ~1h expiry entirely",
], code="""\
// edge/api-proxy.js — new route, existing /api/firefly/generate is untouched
async function handleFireflyGenerateImage(request, env, cors) {
  const { prompt, size } = parsePromptAndSize(await request.json());
  const url = await fireflyGenerate(prompt, size, env);   // shared helper

  const imageRes = await fetch(url);
  return new Response(imageRes.body, {
    status: 200,
    headers: { 'Content-Type': imageRes.headers.get('Content-Type') || 'image/png', ...cors },
  });
}
// Route: POST /api/firefly/generate-image""")

# ── slide 7 — step 3: detecting author context ───────────────────────────────
add_bullet_slide("Step 3 — Detecting Universal Editor Context", [
    "No hostname guessing needed — Universal Editor instruments every block",
    "wrapper with a data-aue-resource attribute, only present inside its canvas",
    "Confirmed live: inspecting the block in Universal Editor showed",
    "data-resource=\"urn:aemconnection:/content/myeds-xwalk/index/",
    "               jcr:content/root/section_158074003/firefly\"",
    "Stripping the urn:aemconnection: prefix gives the exact JCR path",
    "of the block's own content node — this is what later gets patched",
], code="""\
const resource = block.getAttribute('data-aue-resource');

// Published page: attribute absent entirely — never call Firefly there,
// only ever render an already-persisted <picture>.
if (!resource) return;

// Editor canvas: strip the urn prefix to get the real Sling resource path
const resourcePath = resource.replace('urn:aemconnection:', '');
// → /content/myeds-xwalk/index/jcr:content/root/section_158074003/firefly""")

# ── slide 8 — step 4: auto-generate ──────────────────────────────────────────
add_bullet_slide("Step 4 — Auto-Generate on Prompt Entry", [
    "No \"Generate\" button — the first version had one, but felt clunky",
    "for a \"type a prompt, see the image\" authoring flow",
    "Debounced 1.2s after the last field edit — Universal Editor re-runs",
    "decorate() on every properties-rail change, so without debouncing a",
    "single prompt would trigger many redundant Firefly calls mid-typing",
    "An in-flight Set guards against overlapping generations if the author",
    "keeps editing while a generation is already running",
], code="""\
const pendingTimers = new Map();   // resource -> timeoutId
const inFlight = new Set();
const DEBOUNCE_MS = 1200;

if (pendingTimers.has(resource)) clearTimeout(pendingTimers.get(resource));
if (inFlight.has(resource)) return;

const timerId = setTimeout(async () => {
  inFlight.add(resource);
  await generateAndPersist(preview, prompt, aspectRatio, resource);
  window.location.reload();
}, DEBOUNCE_MS);
pendingTimers.set(resource, timerId);""")

# ── slide 9 — step 5: DAM upload ─────────────────────────────────────────────
add_bullet_slide("Step 5 — Upload to AEM DAM", [
    "AEMaaCS deprecated direct binary POST — the supported flow is a",
    "3-step Direct Binary Upload: initiateUpload → PUT to blob storage",
    "→ completeUpload",
    "Same-origin fetch with credentials: 'include' — the author's own",
    "session cookie authenticates every call, no separate service token",
    "Every mutating call needs a Granite CSRF-Token header",
    "Target folder auto-creates itself on first use if missing",
], code="""\
const csrfToken = (await fetch('/libs/granite/csrf/token.json',
  { credentials: 'include' }).then(r => r.json())).token;

let initRes = await initiate();                 // POST <folder>.initiateUpload.json
if (initRes.status === 404) {                    // folder doesn't exist yet
  await createDamFolder(csrfToken);              // POST /api/assets/<folder>
  initRes = await initiate();
}
const { uploadURIs, uploadToken } = (await initRes.json()).files[0];
await fetch(uploadURIs[0], { method: 'PUT', body: blob });   // no auth needed
await fetch(completeURI, { method: 'POST', credentials: 'include',
  headers: { 'CSRF-Token': csrfToken },
  body: new URLSearchParams({ fileName, mimeType, uploadToken }) });""")

# ── slide 10 — step 6: patch content node ────────────────────────────────────
add_bullet_slide("Step 6 — Patch the Block's Own Content Node", [
    "Standard Sling Default POST Servlet behaviour: POST form fields named",
    "after JCR properties updates them on that resource",
    "Same CSRF token and same-origin credentials as the DAM upload",
    "This is additive — it only ever adds/updates the \"image\" property,",
    "never touches the author's prompt/aspectRatio values",
], code="""\
await fetch(resourcePath, {
  method: 'POST',
  credentials: 'include',
  headers: { 'Content-Type': 'application/x-www-form-urlencoded',
             'CSRF-Token': csrfToken },
  body: new URLSearchParams({ image: assetPath }),
});
// resourcePath = /content/myeds-xwalk/index/jcr:content/root/section_.../firefly
// assetPath    = /content/dam/myeds-xwalk/generated/firefly-<timestamp>.jpg""")

# ── slide 11 — bugs table ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, WHITE)
add_rect(s, 0, 0, Inches(0.18), H, RED)
add_text(s, "Bugs Hit Along the Way",
         Inches(0.35), Inches(0.3), Inches(12.5), Inches(0.8),
         size=28, bold=True, color=RED)
sep = s.shapes.add_shape(1, Inches(0.35), Inches(1.05), Inches(12.6), Pt(2))
sep.fill.solid()
sep.fill.fore_color.rgb = RED
sep.line.fill.background()

headers = ["Symptom", "Root Cause", "Fix"]
rows_data = [
    ["Broken image icon shown\nbefore any generation ran",
     ".ai-image-preview's own\ndisplay:block overrode the\nhidden attribute",
     "Add .ai-image-preview[hidden]\n{ display: none; }"],
    ["Generate button rendered\nwith invisible text",
     "Global button{} rule in\nstyles.css sets color:\nvar(--background-color)",
     "Explicitly set color on\n.ai-image-generate\n(later: button removed entirely)"],
    ["\"Something went wrong…\"\nwith no further detail",
     "Catch block swallowed the\nreal error into one generic\nmessage",
     "Surface e.message directly\nin the status text + console.error"],
    ["\"Failed to fetch\" calling\nthe Worker from the editor",
     "Worker's CORS allowlist\ndidn't include the AEM\nauthor origin",
     "Add author-p42808-...\n.adobeaemcloud.com to\nALLOWED_ORIGINS"],
]

col_w = [Inches(4.0), Inches(4.2), Inches(4.2)]
col_x = [Inches(0.35), Inches(4.45), Inches(8.75)]
row_h = Inches(1.15)
header_y = Inches(1.2)

for hdr, cx, cw in zip(headers, col_x, col_w):
    add_rect(s, cx, header_y, cw - Inches(0.05), Inches(0.45), BLACK)
    add_text(s, hdr, cx + Inches(0.08), header_y + Inches(0.05),
             cw - Inches(0.15), Inches(0.38),
             size=14, bold=True, color=WHITE)

for ri, row in enumerate(rows_data):
    ry = header_y + Inches(0.45) + ri * row_h
    bg = LGREY if ri % 2 == 0 else WHITE
    for cell, cx, cw in zip(row, col_x, col_w):
        add_rect(s, cx, ry, cw - Inches(0.05), row_h - Inches(0.05), bg)
        add_text(s, cell, cx + Inches(0.08), ry + Inches(0.06),
                 cw - Inches(0.18), row_h - Inches(0.1),
                 size=11.5, color=DGREY)

add_text(s, "Adobe", Inches(0.35), H - Inches(0.5), Inches(1.5), Inches(0.4),
         size=13, bold=True, color=RED)
add_text(s, "©2026 Adobe. All Rights Reserved. Adobe Confidential.",
         Inches(8), H - Inches(0.5), Inches(5), Inches(0.4),
         size=9, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)

# ── slide 12 — current status / open issue ───────────────────────────────────
add_bullet_slide("Current Status & Open Issue", [
    "Confirmed working: prompt entry, debounced auto-generate, inline",
    "preview, CSRF token fetch, and the DAM upload's initiateUpload step",
    "",
    "Open issue: after completeUpload succeeds, AEM's asset-processing",
    "workflow (thumbnail/rendition generation) starts and has been",
    "observed to fail — the image is then no longer displayed after reload",
    "",
    "This is server-side AEM workflow behaviour, not something visible from",
    "the block's own code — next step is checking Tools → Operations →",
    "Workflow → Failures (or the Assets UI directly) for the actual",
    "processing error, since the folder was auto-created by our code and",
    "may not inherit the same ACLs/workflow launcher config as a normal",
    "authored DAM folder",
])

# ── slide 13 — best practices summary table ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, WHITE)
add_rect(s, 0, 0, Inches(0.18), H, RED)
add_text(s, "Best Practices & Summary",
         Inches(0.35), Inches(0.3), Inches(12.5), Inches(0.8),
         size=28, bold=True, color=RED)
sep = s.shapes.add_shape(1, Inches(0.35), Inches(1.05), Inches(12.6), Pt(2))
sep.fill.solid()
sep.fill.fore_color.rgb = RED
sep.line.fill.background()

headers = ["Step", "Action", "Notes"]
rows_data = [
    ["1. Model",        "prompt, aspectRatio,\nimage (reference)",              "image is written\nprogrammatically, not picked"],
    ["2. Worker route", "Stream bytes instead of\na presigned URL",              "Reuse existing IMS auth —\ndon't duplicate it"],
    ["3. Author context", "Check data-aue-resource\non the block",              "Absent on published pages\n— never call Firefly there"],
    ["4. Auto-generate", "Debounce + in-flight guard",                          "Prevents redundant calls\nwhile the author is editing"],
    ["5. DAM upload",    "initiateUpload / PUT /\ncompleteUpload + CSRF",        "Same-origin, cookie auth —\nno extra service token"],
    ["6. Persist",       "Sling POST the block's\nown resource path",           "Additive — only touches\nthe image property"],
    ["7. Debug",         "Surface real error text,\nnot a generic message",     "CORS/CSRF/workflow failures\nall look identical otherwise"],
]

col_w = [Inches(2.1), Inches(3.8), Inches(4.5)]
col_x = [Inches(0.35), Inches(2.5), Inches(6.35)]
row_h = Inches(0.68)
header_y = Inches(1.2)

for hdr, cx, cw in zip(headers, col_x, col_w):
    add_rect(s, cx, header_y, cw - Inches(0.05), Inches(0.45), BLACK)
    add_text(s, hdr, cx + Inches(0.08), header_y + Inches(0.05),
             cw - Inches(0.15), Inches(0.38),
             size=14, bold=True, color=WHITE)

for ri, row in enumerate(rows_data):
    ry = header_y + Inches(0.45) + ri * row_h
    bg = LGREY if ri % 2 == 0 else WHITE
    for cell, cx, cw in zip(row, col_x, col_w):
        add_rect(s, cx, ry, cw - Inches(0.05), row_h - Inches(0.05), bg)
        add_text(s, cell, cx + Inches(0.08), ry + Inches(0.06),
                 cw - Inches(0.18), row_h - Inches(0.1),
                 size=11.5, color=DGREY)

add_text(s, "Adobe", Inches(0.35), H - Inches(0.5), Inches(1.5), Inches(0.4),
         size=13, bold=True, color=RED)
add_text(s, "©2026 Adobe. All Rights Reserved. Adobe Confidential.",
         Inches(8), H - Inches(0.5), Inches(5), Inches(0.4),
         size=9, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)

# ── slide 14 — end ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, W, H, RED)
add_text(s, "Adobe",
         Inches(5.9), Inches(3.3), Inches(1.6), Inches(0.7),
         size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── save ──────────────────────────────────────────────────────────────────────
out = "/Users/sevin/Library/CloudStorage/OneDrive-Adobe/NEW/Workspace/EDS/myeds-xwalk/docs/AI-Image-Block-Build-Guide.pptx"
prs.save(out)
print(f"Saved → {out}")
